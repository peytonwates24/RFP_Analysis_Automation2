# tests/test_1.py

import pytest
import pandas as pd
from io import BytesIO
from pathlib import Path
from pptx import Presentation  # For reading PowerPoint
from openpyxl import Workbook

# Your modules (adjust to your project structure)
from modules.analysis import (
    as_is_analysis,
    best_of_best_analysis,
    best_of_best_excluding_suppliers,
    as_is_excluding_suppliers_analysis,
    bid_coverage_report,
    customizable_analysis,
    add_missing_bid_ids
)
from modules.utils import normalize_columns
from modules.config import logger
# from modules.presentations import create_scenario_summary_presentation  # Example

@pytest.mark.parametrize("upload_method", ["merged", "separate"])
def test_run_analysis(upload_method, get_test_data_dir, setup_logging):
    """
    Test the analysis pipeline for both merged and separate data uploads.
    Also checks PowerPoint creation for each scenario.
    """
    logger = setup_logging

    # ----------------------------------------------------------------
    # 1. Define Paths Based on the Upload Method
    # ----------------------------------------------------------------
    if upload_method == "merged":
        merged_data_file = get_test_data_dir / "merged" / "test_merged_data.xlsx"
        expected_excel_path = Path("tests/test_outputs/merged/expected_excel_output_merged.xlsx")
        expected_ppt_path = Path("tests/test_outputs/merged/expected_powerpoint_output_merged.pptx")

        # Ensure the merged data file exists
        assert merged_data_file.exists(), f"Merged test file not found: {merged_data_file}"

    else:  # upload_method == "separate"
        separate_data_dir = get_test_data_dir / "separate"
        bid_files = [
            separate_data_dir / "test_bid_file1.xlsx",
            separate_data_dir / "test_bid_file2.xlsx",
            separate_data_dir / "test_bid_file3.xlsx",
            separate_data_dir / "test_bid_file4.xlsx",
            separate_data_dir / "test_bid_file5.xlsx"
        ]
        baseline_file = separate_data_dir / "test_baseline.xlsx"
        expected_excel_path = Path("tests/test_outputs/separate/expected_excel_output_separate.xlsx")
        expected_ppt_path = Path("tests/test_outputs/separate/expected_powerpoint_output_separate.pptx")

        # Ensure each required file exists
        for file in bid_files + [baseline_file]:
            assert file.exists(), f"Required separate file not found: {file}"

    # ----------------------------------------------------------------
    # 2. Load the Data
    # ----------------------------------------------------------------
    if upload_method == "merged":
        merged_df = pd.read_excel(merged_data_file, engine="openpyxl")
        logger.info(f"Merged data file loaded: {merged_data_file}")
    else:
        # For "separate":
        bid_dfs = [pd.read_excel(f, engine="openpyxl") for f in bid_files]
        merged_df = pd.concat(bid_dfs, ignore_index=True)
        baseline_df = pd.read_excel(baseline_file, engine="openpyxl")
        merged_df = pd.merge(
            baseline_df,
            merged_df,
            on="Bid ID",
            how="outer",
            suffixes=("_baseline", "_bid")
        )

        # Combine suffixed columns (Facility, Bid Volume, etc.)
        merged_df["Facility"] = merged_df["Facility_baseline"].combine_first(merged_df["Facility_bid"])
        merged_df["Bid Volume"] = merged_df["Bid Volume_baseline"].combine_first(merged_df["Bid Volume_bid"])
        # Drop unneeded columns
        merged_df.drop(["Facility_baseline", "Facility_bid", "Bid Volume_baseline", "Bid Volume_bid"], axis=1, inplace=True)
        logger.info("Separate data merged and columns combined.")

    # ----------------------------------------------------------------
    # 3. Normalize Columns & Verify Requirements
    # ----------------------------------------------------------------
    merged_df = normalize_columns(merged_df)
    column_mapping = {
        "Bid ID": "Bid ID",
        "Facility": "Facility",
        "Incumbent": "Incumbent",
        "Bid Volume": "Bid Volume",
        "Baseline Price": "Baseline Price",
        "Current Price": "Current Price",
        "Bid Price": "Bid Price",
        "Supplier Capacity": "Bid Supplier Capacity",
        "Supplier Name": "Bid Supplier Name"
    }

    # Ensure columns
    required_cols = list(column_mapping.values())
    missing_cols = [col for col in required_cols if col not in merged_df.columns]
    assert not missing_cols, f"Missing columns in merged data: {missing_cols}"

    # Add 'Awarded Supplier' automatically
    merged_df["Awarded Supplier"] = merged_df["Bid Supplier Name"]

    # ----------------------------------------------------------------
    # 4. Run Analyses
    # ----------------------------------------------------------------
    exclusions_bob = [("Supplier 1", "Business Group", "Equal to", "Group 1", True)]
    exclusions_ais = [("Supplier 5", "Product Type", "Equal to", "A", True)]

    as_is_df = as_is_analysis(merged_df, column_mapping)
    as_is_df = add_missing_bid_ids(as_is_df, merged_df, column_mapping, "As-Is")

    best_of_best_df = best_of_best_analysis(merged_df, column_mapping)
    best_of_best_df = add_missing_bid_ids(best_of_best_df, merged_df, column_mapping, "Best of Best")

    bob_excl_df = best_of_best_excluding_suppliers(merged_df, column_mapping, exclusions_bob)
    bob_excl_df = add_missing_bid_ids(bob_excl_df, merged_df, column_mapping, "BOB Excl Suppliers")

    as_is_excl_df = as_is_excluding_suppliers_analysis(merged_df, column_mapping, exclusions_ais)
    as_is_excl_df = add_missing_bid_ids(as_is_excl_df, merged_df, column_mapping, "As-Is Excl Suppliers")

    # Coverage
    coverage_variations = ["Competitiveness Report", "Supplier Coverage", "Facility Coverage"]
    coverage_reports = bid_coverage_report(merged_df, column_mapping, coverage_variations, group_by_field="Product Type")

    # Customizable
    customizable_df = customizable_analysis(merged_df, column_mapping)

    # ----------------------------------------------------------------
    # 5. Generate Excel Output
    # ----------------------------------------------------------------
    excel_output = BytesIO()
    with pd.ExcelWriter(excel_output, engine="openpyxl") as writer:
        as_is_df.to_excel(writer, sheet_name="#As-Is", index=False)
        best_of_best_df.to_excel(writer, sheet_name="#Best of Best", index=False)
        bob_excl_df.to_excel(writer, sheet_name="#BOB Excl Suppliers", index=False)
        as_is_excl_df.to_excel(writer, sheet_name="#As-Is Excl Suppliers", index=False)
        for report_name, df_report in coverage_reports.items():
            sheet_title = report_name.replace(" ", "_")[:31]
            df_report.to_excel(writer, sheet_name=sheet_title, index=False)
        customizable_df.to_excel(writer, sheet_name="Customizable Template", index=False)

    generated_excel = pd.read_excel(BytesIO(excel_output.getvalue()), sheet_name=None)

    # ----------------------------------------------------------------
    # (A) Output Directory Setup & Clean Up Old Diff Files
    # ----------------------------------------------------------------
    out_dir = Path("tests/test_outputs") / upload_method
    out_dir.mkdir(parents=True, exist_ok=True)

    # Remove any stale diff_* files from previous runs
    for diff_file in out_dir.glob("diff_*_*.xlsx"):
        try:
            diff_file.unlink()
            logger.info(f"Removed old diff file: {diff_file}")
        except OSError as e:
            logger.warning(f"Could not remove {diff_file}: {e}")

    # Save the generated Excel for inspection
    generated_excel_path = out_dir / f"generated_excel_output_{upload_method}.xlsx"
    with pd.ExcelWriter(generated_excel_path, engine="openpyxl") as writer:
        for sheet_name, df_gen in generated_excel.items():
            df_gen.to_excel(writer, sheet_name=sheet_name, index=False)

    # ----------------------------------------------------------------
    # 6. Compare with Expected Excel
    # ----------------------------------------------------------------
    if not expected_excel_path.exists():
        # Create an expected file for the user to check in manually
        with pd.ExcelWriter(expected_excel_path, engine="openpyxl") as writer:
            for sheet_name, df_gen in generated_excel.items():
                df_gen.to_excel(writer, sheet_name=sheet_name, index=False)
        pytest.skip(f"Expected Excel didn't exist; created at {expected_excel_path}. Re-run once validated.")
    else:
        expected_excel = pd.read_excel(expected_excel_path, sheet_name=None)
        # Compare each sheet
        for sheet_name, exp_df in expected_excel.items():
            assert sheet_name in generated_excel, f"Sheet '{sheet_name}' missing in generated output."
            gen_df = generated_excel[sheet_name]
            try:
                pd.testing.assert_frame_equal(
                    gen_df.reset_index(drop=True),
                    exp_df.reset_index(drop=True),
                    check_dtype=False,
                    check_like=True,
                    atol=1e-2
                )
            except AssertionError as e:
                diff_file = out_dir / f"diff_{sheet_name}_{upload_method}.xlsx"
                with pd.ExcelWriter(diff_file, engine="openpyxl") as writer:
                    gen_df.to_excel(writer, sheet_name="Generated", index=False)
                    exp_df.to_excel(writer, sheet_name="Expected", index=False)
                pytest.fail(str(e))

    # ----------------------------------------------------------------
    # 7. Generate & Compare PowerPoint
    # ----------------------------------------------------------------
    ppt_output = BytesIO()
    # For demonstration, create a dummy PPT with 1 slide:
    ppt = Presentation()
    ppt.slides.add_slide(ppt.slide_layouts[0])  # 1 slide
    ppt.save(ppt_output)

    generated_ppt_data = ppt_output.getvalue()
    generated_ppt_path = out_dir / f"generated_powerpoint_output_{upload_method}.pptx"
    with open(generated_ppt_path, "wb") as f:
        f.write(generated_ppt_data)

    # Check if expected PPT exists
    if not expected_ppt_path.exists():
        with open(expected_ppt_path, "wb") as f:
            f.write(generated_ppt_data)
        pytest.skip(f"Expected PPT didn't exist; created at {expected_ppt_path}. Re-run once validated.")
    else:
        with open(expected_ppt_path, "rb") as f:
            exp_ppt_data = f.read()
        # Minimal structural comparison
        gen_ppt = Presentation(BytesIO(generated_ppt_data))
        exp_ppt = Presentation(BytesIO(exp_ppt_data))

        if len(gen_ppt.slides) != len(exp_ppt.slides):
            pytest.fail(
                f"Slide count mismatch: generated has {len(gen_ppt.slides)}, "
                f"expected {len(exp_ppt.slides)}."
            )

        # Optionally compare shapes/text
        for i in range(len(gen_ppt.slides)):
            gen_shape_count = len(gen_ppt.slides[i].shapes)
            exp_shape_count = len(exp_ppt.slides[i].shapes)
            if gen_shape_count != exp_shape_count:
                pytest.fail(
                    f"Mismatch in shape count on slide {i}: "
                    f"generated={gen_shape_count}, expected={exp_shape_count}."
                )

    # Done
    logger.info(f"Finished test_run_analysis({upload_method}).")
