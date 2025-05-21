import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

def format_currency(amount):
    """Formats the currency value with commas."""
    return "${:,.0f}".format(amount)

def format_currency_in_millions(amount):
    """Formats the amount in millions with one decimal and 'MM' suffix."""
    m = amount / 1_000_000
    if amount < 0:
        return f"(${abs(m):,.1f}MM)"
    return f"${m:,.1f}MM"

def add_header(slide, slide_num, title_suffix=""):
    """Add the big blue 'Scenario Summary' header (optionally with a suffix)."""
    left, top, w, h = Inches(0), Inches(0.1), Inches(12.12), Inches(0.42)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if slide_num == 1 and not title_suffix:
        p.text = "Scenario Summary"
    else:
        p.text = f"Scenario Summary{title_suffix or f' #{slide_num}'}"
    p.font.name = "Calibri"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 153)
    tb.fill.background()
    tb.line.fill.background()

def add_row_labels(slide):
    """Add the fixed row labels down the left side."""
    labels = [
        ("Scenario Name",    Inches(0.8)),
        ("Description",      Inches(1.29)),
        ("# of Suppliers (% spend)", Inches(1.9)),
        ("RFP Savings %",    Inches(3.25)),
        ("Total Value Opportunity ($MM)", Inches(3.84)),
        ("Key Considerations", Inches(4.89)),
    ]

    for text, top in labels:
        tb = slide.shapes.add_textbox(Inches(0.3), top, Inches(2.0), Inches(0.3))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 153)
        tb.fill.background()
        tb.line.fill.background()

def process_scenario_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """
    Ensure the optimization-output DataFrame has the columns we need,
    and rename 'Awarded Supplier' -> 'Awarded Supplier'.
    """
    df = df.copy()
    df.columns = df.columns.str.strip()
    df = df.rename(columns={
        "Awarded Supplier":        "Awarded Supplier",
        "Baseline Savings":        "Baseline Savings",
        "Current Price Savings":   "Current Price Savings",
        "Baseline Spend":          "Baseline Spend",
        "Current Baseline Spend":  "Current Baseline Spend",
    })
    for col in [
        "Awarded Supplier",
        "Awarded Supplier Spend",
        "Baseline Savings",
        "Current Price Savings",
        "Baseline Spend",
        "Current Baseline Spend"
    ]:
        if col not in df.columns:
            df[col] = 0
    return df

def add_scenario_content(slide, df, scenario, scenario_position):
    """Adds the content for a single scenario to the slide. scenario_position: 1, 2, or 3"""
    df.columns = df.columns.str.strip()

    expected_columns = [
        'Awarded Supplier',
        'Awarded Supplier Spend',
        'Baseline Savings',
        'Current Price Savings',
        'Baseline Spend',
        'Current Baseline Spend'
    ]
    for col in expected_columns:
        if col not in df.columns:
            raise ValueError(f"Column '{col}' not found in sheet '{scenario}'.")

    base_left = Inches(2.72) + (scenario_position - 1) * Inches(3.15)
    positions = {
        'scenario_name': {'left': base_left, 'top': Inches(0.8), 'width': Inches(2.5), 'height': Inches(0.34)},
        'description': {'left': base_left, 'top': Inches(1.29), 'width': Inches(2.5), 'height': Inches(0.34)},
        'suppliers_entry': {'left': base_left, 'top': Inches(1.85), 'width': Inches(2.5), 'height': Inches(1.11)},
        'rfp_entry': {'left': base_left, 'top': Inches(3.26), 'width': Inches(2.5), 'height': Inches(0.3)},
        'key_considerations': {'left': base_left, 'top': Inches(4.88), 'width': Inches(2.5), 'height': Inches(2.0)},
    }

    # Add Scenario Name
    scenario_name_box = slide.shapes.add_textbox(
        positions['scenario_name']['left'],
        positions['scenario_name']['top'],
        positions['scenario_name']['width'],
        positions['scenario_name']['height']
    )
    scenario_name_tf = scenario_name_box.text_frame
    scenario_name_tf.vertical_anchor = MSO_ANCHOR.TOP
    scenario_name_tf.word_wrap = True
    p = scenario_name_tf.paragraphs[0]
    p.text = scenario
    p.font.size = Pt(14)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 153)  # Dark Blue
    scenario_name_box.fill.background()
    scenario_name_box.line.fill.background()

    # Add Description Entry
    desc_entry = slide.shapes.add_textbox(
        positions['description']['left'],
        positions['description']['top'],
        positions['description']['width'],
        positions['description']['height']
    )
    desc_entry_tf = desc_entry.text_frame
    desc_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    desc_entry_tf.word_wrap = True
    p = desc_entry_tf.paragraphs[0]
    p.text = "Describe your scenario here"
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    desc_entry.fill.background()  # No Fill
    desc_entry.line.fill.background()  # No Line

    # Add Suppliers Entry
    suppliers_entry = slide.shapes.add_textbox(
        positions['suppliers_entry']['left'],
        positions['suppliers_entry']['top'],
        positions['suppliers_entry']['width'],
        positions['suppliers_entry']['height']
    )
    suppliers_entry_tf = suppliers_entry.text_frame
    suppliers_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    suppliers_entry_tf.word_wrap = True
    p = suppliers_entry_tf.paragraphs[0]
    # Calculate number of suppliers and % spend
    total_spend = df['Awarded Supplier Spend'].sum()
    suppliers = df.groupby('Awarded Supplier')['Awarded Supplier Spend'].sum().reset_index()
    suppliers['Spend %'] = suppliers['Awarded Supplier Spend'] / total_spend * 100
    # Sort suppliers by 'Spend %' in descending order
    suppliers = suppliers.sort_values(by='Spend %', ascending=False)
    num_suppliers = suppliers['Awarded Supplier'].nunique()
    supplier_list = [f"{row['Awarded Supplier']} ({row['Spend %']:.0f}%)" for idx, row in suppliers.iterrows()]
    # Print # of suppliers in scenario suppliers joined with % of spend
    supplier_text = f"{num_suppliers}- " + ", ".join(supplier_list)
    p.text = supplier_text
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    suppliers_entry.fill.background()  # No Fill
    suppliers_entry.line.fill.background()  # No Line

    # Calculate AST Savings % and Current Savings %
    total_ast_savings = df['Baseline Savings'].sum()
    total_current_savings = df['Current Price Savings'].sum()
    ast_baseline_spend = df['Baseline Spend'].sum()
    current_baseline_spend = df['Current Baseline Spend'].sum()

    ast_savings_pct = (total_ast_savings / ast_baseline_spend * 100) if ast_baseline_spend != 0 else 0
    current_savings_pct = (total_current_savings / current_baseline_spend * 100) if current_baseline_spend != 0 else 0

    rfp_savings_str = f"{ast_savings_pct:.0f}% | {current_savings_pct:.0f}%"

    # Add RFP Savings % Entry
    rfp_entry = slide.shapes.add_textbox(
        positions['rfp_entry']['left'],
        positions['rfp_entry']['top'],
        positions['rfp_entry']['width'],
        positions['rfp_entry']['height']
    )
    rfp_entry_tf = rfp_entry.text_frame
    rfp_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    rfp_entry_tf.word_wrap = True
    p = rfp_entry_tf.paragraphs[0]
    p.text = rfp_savings_str
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    rfp_entry.fill.background()  # No Fill
    rfp_entry.line.fill.background()  # No Line

    # Add Key Considerations Entry
    key_entry = slide.shapes.add_textbox(
        positions['key_considerations']['left'],
        positions['key_considerations']['top'],
        positions['key_considerations']['width'],
        positions['key_considerations']['height']
    )
    key_entry_tf = key_entry.text_frame
    key_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    key_entry_tf.word_wrap = True
    p = key_entry_tf.paragraphs[0]
    p.text = f"Key considerations for {scenario}"
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    key_entry.fill.background()  # No Fill
    key_entry.line.fill.background()  # No Line

def add_chart(slide, scenario_names, ast_savings_list, current_savings_list):
    """Adds a clustered column chart between RFP Savings % and Key Considerations."""
    chart_data = CategoryChartData()
    chart_data.categories = scenario_names
    chart_data.add_series('AST',     ast_savings_list)
    chart_data.add_series('Current', current_savings_list)

    x, y, cx, cy = Inches(1.86), Inches(3.69), Inches(9.0), Inches(1.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # color the two series
    s1 = chart.series[0]
    s1.format.fill.solid()
    s1.format.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    s2 = chart.series[1]
    s2.format.fill.solid()
    s2.format.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

    # legend on the right
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = "Calibri"

    # no title
    chart.has_title = False

    # hide both axes
    chart.category_axis.visible = False
    chart.value_axis.visible    = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(192,192,192)

    # overlap & gap width
    plot = chart.plots[0]
    plot.gap_width = 217
    plot.overlap   = -27

    # data labels in MM
    data_lists = [ast_savings_list, current_savings_list]
    for idx, series in enumerate(chart.series):
        for pt_idx, pt in enumerate(series.points):
            lbl = pt.data_label
            lbl.has_value = False
            lbl.number_format_is_linked = False
            lbl.visible = True

            tf = lbl.text_frame
            tf.clear()  # remove existing runs
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = format_currency_in_millions(data_lists[idx][pt_idx])
            run.font.size = Pt(12)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0,0,0)


def create_scenario_summary_slides(prs: Presentation, scenario_dfs: dict) -> Presentation:
    """
    Take an existing Presentation (template or blank),
    lay out up to 3 scenarios per slide, repeating until done.
    """
    names = list(scenario_dfs.keys())
    per_slide = 3
    total_slides = (len(names)+per_slide-1)//per_slide

    idx = 0
    for slide_num in range(total_slides):
        if slide_num == 0 and prs.slides:
            slide = prs.slides[0]
            # clear existing shapes
            for shp in list(slide.shapes):
                slide.shapes._spTree.remove(shp._element)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

        add_header(slide, slide_num+1)
        add_row_labels(slide)

        batch_names = []
        base_vals  = []
        curr_vals  = []

        for pos in range(1, per_slide+1):
            if idx >= len(names):
                break
            nm = names[idx]
            df = scenario_dfs[nm]
            df = process_scenario_dataframe(df)

            add_scenario_content(slide, df, nm, pos)

            base_vals.append(df["Baseline Savings"].sum())
            curr_vals.append(df["Current Price Savings"].sum())
            batch_names.append(nm)
            idx += 1

        add_chart(slide, batch_names, base_vals, curr_vals)

    return prs

def create_scenario_summary_presentation(
    scenario_dfs: dict,
    template_file_path: str = None
) -> Presentation:
    """
    Load a template if found, else blank.  Then build all summary slides.
    """
    if template_file_path and os.path.exists(template_file_path):
        prs = Presentation(template_file_path)
    else:
        prs = Presentation()
    return create_scenario_summary_slides(prs, scenario_dfs)


# ppt_scenario.py

def add_scenario_detail_slides(
    prs: Presentation,
    scenario_dfs: dict,
    grouping_col: str,
    template_slide_layout_index: int,
    detail_columns: list = [],
    item_attr_df: pd.DataFrame = None,
    bid_data_df: pd.DataFrame = None
):
    import pandas as pd
    from itertools import cycle
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.chart.data import ChartData
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor

    # 1) Build a consistent color map for all suppliers/incumbents across slides
    all_names = set()
    for raw_df in scenario_dfs.values():
        tmp = raw_df.copy()
        tmp.columns = tmp.columns.str.strip()
        if "Awarded Supplier Name" in tmp:
            tmp.rename(columns={"Awarded Supplier Name": "Awarded Supplier"}, inplace=True)
        all_names.update(tmp["Awarded Supplier"].dropna().unique())
        if "Incumbent" in tmp:
            all_names.update(tmp["Incumbent"].dropna().unique())
    palette = cycle([
        RGBColor(68,114,196),
        RGBColor(237,125,49),
        RGBColor(165,165,165),
        RGBColor(255,192,0),
        RGBColor(112,173,71),
    ])
    supplier_color_map = {name: next(palette) for name in sorted(all_names)}

    # 2) Find the "Default" layout by name, fallback to provided index
    try:
        default_layout = next(
            l for l in prs.slide_layouts
            if l.name.lower() == "default"
        )
    except StopIteration:
        default_layout = prs.slide_layouts[template_slide_layout_index]

    # 3) Iterate scenarios
    for scen_name, raw_df in scenario_dfs.items():
        df = raw_df.copy()
        df.columns = df.columns.str.strip()

        # 4) Merge in grouping & detail columns if missing
        for col in [grouping_col] + detail_columns:
            if col not in df.columns:
                # from Item Attributes by Bid ID
                if item_attr_df is not None and col in item_attr_df.columns:
                    df = df.merge(
                        item_attr_df[['Bid ID', col]].drop_duplicates(),
                        on='Bid ID', how='left'
                    )
                # from Bid Data by Bid ID + Awarded Supplier
                elif bid_data_df is not None and col in bid_data_df.columns:
                    df = df.merge(
                        bid_data_df[['Bid ID', 'Supplier Name', col]].drop_duplicates(),
                        left_on=['Bid ID', 'Awarded Supplier'],
                        right_on=['Bid ID', 'Supplier Name'],
                        how='left'
                    ).drop(columns=['Supplier Name'])
                else:
                    df[col] = ""

        # ensure grouping & detail columns exist now
        df[grouping_col] = df.get(grouping_col, "")
        for col in detail_columns:
            df[col] = df.get(col, "")

        # normalize key columns
        if "Awarded Supplier Name" in df:
            df.rename(columns={"Awarded Supplier Name": "Awarded Supplier"}, inplace=True)
        if "Bid Volume" in df and "Awarded Volume" not in df:
            df["Awarded Volume"] = df["Bid Volume"]
        if "Current Price Savings" in df and "Current Savings" not in df:
            df.rename(columns={"Current Price Savings": "Current Savings"}, inplace=True)
        if "Discounted Awarded Supplier Price" in df and "Effective Supplier Price" not in df:
            df.rename(columns={"Discounted Awarded Supplier Price": "Effective Supplier Price"}, inplace=True)

        # add slide and clear its shapes
        slide = prs.slides.add_slide(default_layout)
        for shp in list(slide.shapes):
            slide.shapes._spTree.remove(shp._element)

        # Header
        hdr = slide.shapes.add_textbox(
            Inches(0), Inches(0.01),
            Inches(12.5), Inches(0.5)
        )
        tf = hdr.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = f"{scen_name} Scenario Details"
        p.font.name, p.font.size, p.font.bold = "Calibri", Pt(25), True
        p.font.color.rgb = RGBColor(0,51,153)
        hdr.fill.background(); hdr.line.fill.background()

        # Stacked bar chart (10.5" wide, no title, legend on top, 12pt legend, 8pt labels)
        total_vol = df["Awarded Volume"].sum() or 1
        vol_df = df.groupby("Awarded Supplier")["Awarded Volume"].sum().reset_index()
        vol_df["pct"] = vol_df["Awarded Volume"] / total_vol

        cd = ChartData(); cd.categories = ['']
        for _, r in vol_df.iterrows():
            cd.add_series(r["Awarded Supplier"], [r["pct"]])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(0), Inches(0.3),
            Inches(10.5), Inches(1.7),
            cd
        ).chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = "Calibri"

        for series in chart.series:
            clr = supplier_color_map.get(series.name, RGBColor(158,158,158))
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = clr
            series.has_data_labels = True
            dl = series.data_labels
            dl.number_format = '0%'
            dl.position = XL_LABEL_POSITION.INSIDE_BASE
            dl.font.size = Pt(8)
            dl.font.color.rgb = RGBColor(255,255,255)

        # Info boxes at top = 0.65"
        num_items = df["Bid ID"].nunique()
        num_trans = df[df["Awarded Supplier"] != df["Incumbent"]]["Bid ID"].nunique()

        box_i = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.97), Inches(0.65),
            Inches(2.0), Inches(1.0)
        )
        box_i.fill.solid(); box_i.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
        box_i.line.fill.background()
        tb_i = box_i.text_frame; tb_i.text = f"# of items\n{num_items}"
        tb_i.vertical_anchor = MSO_ANCHOR.MIDDLE
        for par in tb_i.paragraphs:
            par.alignment = PP_ALIGN.CENTER
            par.font.size = Pt(20)
            par.font.color.rgb = RGBColor(255,255,255)

        box_t = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(11.15), Inches(0.65),
            Inches(2.0), Inches(1.0)
        )
        box_t.fill.solid(); box_t.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        box_t.line.fill.background()
        tb_t = box_t.text_frame; tb_t.text = f"# of Transitions\n{num_trans}"
        tb_t.vertical_anchor = MSO_ANCHOR.MIDDLE
        for par in tb_t.paragraphs:
            par.alignment = PP_ALIGN.CENTER
            par.font.size = Pt(20)
            par.font.color.rgb = RGBColor(255,255,255)

        # Build summary data
        summary = []
        for grp, sub in df.groupby(grouping_col, sort=False):
            vol = sub["Awarded Volume"].sum()
            avg_cur = sub["Current Price"].mean() if "Current Price" in sub else 0
            cur_sp = avg_cur * vol
            avg_bid = sub["Effective Supplier Price"].mean() if "Effective Supplier Price" in sub else 0
            bid_sp = avg_bid * vol
            sav = sub["Current Savings"].sum()
            summary.append({
                "Grouping": grp,
                **{col: ", ".join(sub[col].dropna().astype(str).unique()) for col in detail_columns},
                "Bid Volume": vol,
                "Avg Current Price": avg_cur,
                "Current Spend": cur_sp,
                "Avg Bid Price": avg_bid,
                "Bid Spend": bid_sp,
                "Current Savings": sav,
                "Incumbent Dist": sub.groupby("Incumbent")["Awarded Volume"].sum().reset_index(),
                "Awarded Dist":   sub.groupby("Awarded Supplier")["Awarded Volume"].sum().reset_index()
            })

        # Table
        cols = ["Grouping"] + detail_columns + [
            "Bid Volume","Incumbent","Avg Current Price","Current Spend",
            "Awarded Supplier","Avg Bid Price","Bid Spend","Current Savings"
        ]
        rows = len(summary) + 2
        left, top = Inches(0.5), Inches(2.0)
        table = slide.shapes.add_table(
            rows, len(cols),
            left, top,
            Inches(12.5), Inches(1.0)
        ).table
        table.columns[2].width = Inches(1.5)
        table.columns[5].width = Inches(1.5)
        table.rows[0].height = Inches(0.38)
        for i in range(1, len(summary)+1):
            table.rows[i].height = Inches(1.0)
        table.rows[len(summary)+1].height = Inches(0.38)

        # Headers
        for j, hdr in enumerate(cols):
            cell = table.cell(0, j)
            cell.text = hdr
            ph = cell.text_frame.paragraphs[0]
            ph.font.bold = True
            ph.font.size = Pt(12)
            ph.font.name = "Calibri"
            cell.alignment = PP_ALIGN.CENTER

        # Populate rows & overlay pies
        for i, data in enumerate(summary, start=1):
            for j, col in enumerate(cols):
                if col not in ("Incumbent","Awarded Supplier"):
                    val = data.get(col, "")
                    if col in ("Avg Current Price","Avg Bid Price"):
                        txt = f"${val:.2f}"
                    elif col in ("Current Spend","Bid Spend","Current Savings"):
                        txt = f"${(val/1_000_000):.1f}MM"
                    else:
                        txt = str(val)
                    cell = table.cell(i, j)
                    cell.text = txt
                    pf = cell.text_frame.paragraphs[0]
                    pf.alignment = PP_ALIGN.CENTER
                    pf.font.size = Pt(10)

            for dist_key, col_idx in [("Incumbent Dist",2),("Awarded Dist",5)]:
                dist = data[dist_key]
                cd2 = ChartData()
                cd2.categories = dist.iloc[:,0].tolist()
                cd2.add_series("", dist.iloc[:,1].tolist())
                cell_left = left + sum(table.columns[x].width for x in range(col_idx))
                cell_top  = top + table.rows[0].height + table.rows[i].height*(i-1)
                cell_w = table.columns[col_idx].width
                cell_h = table.rows[i].height
                pie = slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE,
                    cell_left, cell_top,
                    cell_w, cell_h,
                    cd2
                ).chart
                pie.has_title = pie.has_legend = False
                for idx_pt, point in enumerate(pie.series[0].points):
                    cat = dist.iloc[idx_pt,0]
                    clr = supplier_color_map.get(cat, RGBColor(192,192,192))
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = clr

        # Totals row
        tr = len(summary) + 1
        vol_tot = sum(r["Bid Volume"] for r in summary)
        cur_tot = sum(r["Current Spend"] for r in summary)
        bid_tot = sum(r["Bid Spend"] for r in summary)
        sav_tot = sum(r["Current Savings"] for r in summary)
        avg_cur_tot = sum(r["Avg Current Price"]*r["Bid Volume"] for r in summary)/vol_tot
        avg_bid_tot = sum(r["Avg Bid Price"]*r["Bid Volume"] for r in summary)/vol_tot

        totals = {
            "Grouping": "Totals",
            "Bid Volume": vol_tot,
            "Avg Current Price": avg_cur_tot,
            "Current Spend": cur_tot,
            "Avg Bid Price": avg_bid_tot,
            "Bid Spend": bid_tot,
            "Current Savings": sav_tot
        }
        for j, col in enumerate(cols):
            cell = table.cell(tr, j)
            val = totals.get(col, "")
            if isinstance(val, (int, float)):
                if col in ("Avg Current Price","Avg Bid Price"):
                    text = f"${val:.2f}"
                else:
                    text = f"${(val/1_000_000):.1f}MM"
            else:
                text = str(val)
            cell.text = text
            ph = cell.text_frame.paragraphs[0]
            ph.font.bold = True
            ph.font.size = Pt(12)
            ph.font.name = "Calibri"
            ph.alignment = PP_ALIGN.CENTER