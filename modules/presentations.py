import pandas as pd
from pptx import Presentation
from .config import logger
import streamlit as st
import os
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.workbook.defined_name import DefinedName
from openpyxl.styles import Font
from openpyxl import Workbook
from openpyxl.utils import get_column_letter, column_index_from_string
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import tempfile
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import ChartData, CategoryChartData
from itertools import cycle
from collections import OrderedDict


# //// Scenario Summary Presentation /////#

def print_slide_layouts(template_file_path):
    from pptx import Presentation
    prs = Presentation(template_file_path)
    for index, layout in enumerate(prs.slide_layouts):
        print(f"Index {index}: Layout Name - '{layout.name}'")

def format_currency(amount):
    """Formats the currency value with commas."""
    return "${:,.0f}".format(amount)

def format_currency_in_millions(amount):
    """Formats the amount in millions with one decimal place and appends 'MM'."""
    amount_in_millions = amount / 1_000_000
    if amount < 0:
        return f"(${abs(amount_in_millions):,.1f}MM)"
    else:
        return f"${amount_in_millions:,.1f}MM"

def add_header(slide, slide_num, title_suffix=""):
    """Adds the main header to the slide, optionally including a title suffix for sub-summaries."""
    left = Inches(0)
    top = Inches(0.01)
    width = Inches(12.12)
    height = Inches(0.42)
    header = slide.shapes.add_textbox(left, top, width, height)
    header_tf = header.text_frame
    header_tf.vertical_anchor = MSO_ANCHOR.TOP
    header_tf.word_wrap = True
    p = header_tf.paragraphs[0]
    if slide_num == 1 and not title_suffix:
        p.text = "Scenario Summary"
    else:
        # If title_suffix is provided, incorporate it into the header.
        # Example: "Scenario Summary (Port Hudson)"
        if title_suffix:
            p.text = f"Scenario Summary {title_suffix}"
        else:
            p.text = f"Scenario Summary #{slide_num}"
    p.font.size = Pt(25)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 51, 153)  # Dark Blue
    p.font.bold = True
    header.fill.background()  # No Fill
    header.line.fill.background()  # No Line



def add_row_labels(slide):
    """Adds row labels to the slide (once per slide)."""
    row_labels = [
        'Scenario Name',
        'Description',
        '# of Suppliers (% spend)',
        'RFP Savings %',
        'Total Value Opportunity ($MM)',
        'Key Considerations'
    ]
    top_positions = [
        Inches(0.8),
        Inches(1.29),
        Inches(1.9),
        Inches(3.25),
        Inches(3.84),
        Inches(4.89) 
    ]
    left = Inches(0.58)
    width = Inches(1.5)
    height = Inches(0.2)
    for label_text, top in zip(row_labels, top_positions):
        label_box = slide.shapes.add_textbox(left, top, width, height)
        label_tf = label_box.text_frame
        label_tf.vertical_anchor = MSO_ANCHOR.TOP
        label_tf.word_wrap = True
        p = label_tf.paragraphs[0]
        p.text = label_text
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0, 51, 153)  # Dark Blue
        label_box.fill.background()
        label_box.line.fill.background()

def add_scenario_content(slide, df, scenario, scenario_position):
    """Adds the content for a single scenario to the slide. scenario_position: 1, 2, or 3"""
    df.columns = df.columns.str.strip()

    expected_columns = [
        'Awarded Supplier Name',
        'Awarded Supplier Spend',
        'AST Savings',
        'Current Savings',
        'AST Baseline Spend',
        'Current Baseline Spend'
    ]
    for col in expected_columns:
        if col not in df.columns:
            raise ValueError(f"Column '{col}' not found in sheet '{scenario}'.")

    base_left = Inches(2.72) + (scenario_position - 1) * Inches(3.15)
    positions = {
        'scenario_name': {'left': base_left, 'top': Inches(0.8), 'width': Inches(2.5), 'height': Inches(0.34)},
        'description': {'left': base_left, 'top': Inches(1.29), 'width': Inches(2.5), 'height': Inches(0.34)},
        'suppliers_entry': {'left': base_left, 'top': Inches(1.85), 'width': Inches(2.5), 'height': Inches(1.11)},
        'rfp_entry': {'left': base_left, 'top': Inches(3.26), 'width': Inches(2.5), 'height': Inches(0.3)},
        'key_considerations': {'left': base_left, 'top': Inches(4.88), 'width': Inches(2.5), 'height': Inches(2.0)},
    }

    # Add Scenario Name
    scenario_name_box = slide.shapes.add_textbox(
        positions['scenario_name']['left'],
        positions['scenario_name']['top'],
        positions['scenario_name']['width'],
        positions['scenario_name']['height']
    )
    scenario_name_tf = scenario_name_box.text_frame
    scenario_name_tf.vertical_anchor = MSO_ANCHOR.TOP
    scenario_name_tf.word_wrap = True
    p = scenario_name_tf.paragraphs[0]
    p.text = scenario
    p.font.size = Pt(14)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 153)  # Dark Blue
    scenario_name_box.fill.background()
    scenario_name_box.line.fill.background()

    # Add Description Entry
    desc_entry = slide.shapes.add_textbox(
        positions['description']['left'],
        positions['description']['top'],
        positions['description']['width'],
        positions['description']['height']
    )
    desc_entry_tf = desc_entry.text_frame
    desc_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    desc_entry_tf.word_wrap = True
    p = desc_entry_tf.paragraphs[0]
    p.text = "Describe your scenario here"
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    desc_entry.fill.background()  # No Fill
    desc_entry.line.fill.background()  # No Line

    # Add Suppliers Entry
    suppliers_entry = slide.shapes.add_textbox(
        positions['suppliers_entry']['left'],
        positions['suppliers_entry']['top'],
        positions['suppliers_entry']['width'],
        positions['suppliers_entry']['height']
    )
    suppliers_entry_tf = suppliers_entry.text_frame
    suppliers_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    suppliers_entry_tf.word_wrap = True
    p = suppliers_entry_tf.paragraphs[0]
    # Calculate number of suppliers and % spend
    total_spend = df['Awarded Supplier Spend'].sum()
    suppliers = df.groupby('Awarded Supplier Name')['Awarded Supplier Spend'].sum().reset_index()
    suppliers['Spend %'] = suppliers['Awarded Supplier Spend'] / total_spend * 100
    # Sort suppliers by 'Spend %' in descending order
    suppliers = suppliers.sort_values(by='Spend %', ascending=False)
    num_suppliers = suppliers['Awarded Supplier Name'].nunique()
    supplier_list = [f"{row['Awarded Supplier Name']} ({row['Spend %']:.0f}%)" for idx, row in suppliers.iterrows()]
    # Print # of suppliers in scenario suppliers joined with % of spend
    supplier_text = f"{num_suppliers}- " + ", ".join(supplier_list)
    p.text = supplier_text
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    suppliers_entry.fill.background()  # No Fill
    suppliers_entry.line.fill.background()  # No Line

    # Calculate AST Savings % and Current Savings %
    total_ast_savings = df['AST Savings'].sum()
    total_current_savings = df['Current Savings'].sum()
    ast_baseline_spend = df['AST Baseline Spend'].sum()
    current_baseline_spend = df['Current Baseline Spend'].sum()

    ast_savings_pct = (total_ast_savings / ast_baseline_spend * 100) if ast_baseline_spend != 0 else 0
    current_savings_pct = (total_current_savings / current_baseline_spend * 100) if current_baseline_spend != 0 else 0

    rfp_savings_str = f"{ast_savings_pct:.0f}% | {current_savings_pct:.0f}%"

    # Add RFP Savings % Entry
    rfp_entry = slide.shapes.add_textbox(
        positions['rfp_entry']['left'],
        positions['rfp_entry']['top'],
        positions['rfp_entry']['width'],
        positions['rfp_entry']['height']
    )
    rfp_entry_tf = rfp_entry.text_frame
    rfp_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    rfp_entry_tf.word_wrap = True
    p = rfp_entry_tf.paragraphs[0]
    p.text = rfp_savings_str
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    rfp_entry.fill.background()  # No Fill
    rfp_entry.line.fill.background()  # No Line

    # Add Key Considerations Entry
    key_entry = slide.shapes.add_textbox(
        positions['key_considerations']['left'],
        positions['key_considerations']['top'],
        positions['key_considerations']['width'],
        positions['key_considerations']['height']
    )
    key_entry_tf = key_entry.text_frame
    key_entry_tf.vertical_anchor = MSO_ANCHOR.TOP
    key_entry_tf.word_wrap = True
    p = key_entry_tf.paragraphs[0]
    p.text = f"Key considerations for {scenario}"
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    key_entry.fill.background()  # No Fill
    key_entry.line.fill.background()  # No Line

def add_chart(slide, scenario_names, ast_savings_list, current_savings_list):
    """Adds a bar chart to the slide between RFP Savings % and Key Considerations."""
    chart_data = CategoryChartData()
    chart_data.categories = scenario_names

    # Use raw values for chart data
    chart_data.add_series('AST', ast_savings_list)
    chart_data.add_series('Current', current_savings_list)

    # Define the chart size and position
    x, y, cx, cy = Inches(1.86), Inches(3.69), Inches(9.0), Inches(1.0)  # Adjusted position and size

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Format the chart
    # Set the colors for the series
    # AST Savings in Dark Blue Accent 1
    # Current Savings in Turquoise Accent 2

    # For AST Savings (Series 1)
    series1 = chart.series[0]
    fill1 = series1.format.fill
    fill1.solid()
    fill1.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1  # Dark Blue Accent 1

    # For Current Savings (Series 2)
    series2 = chart.series[1]
    fill2 = series2.format.fill
    fill2.solid()
    fill2.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2  # Turquoise Accent 2

    # Set the legend position
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False  # "Show the legend without overlapping the chart"
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = "Calibri"

    # Remove chart title
    chart.has_title = False

    # Hide the horizontal (category) axis
    category_axis = chart.category_axis
    category_axis.visible = False  # Hide the horizontal axis completely

    # Configure the vertical (value) axis
    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(192, 192, 192)  # Light Gray
    value_axis.tick_labels.visible = False  # Hide labels

    # Adjust series overlap and gap width
    plot = chart.plots[0]
    plot.gap_width = 217  # Adjust gap width to 217%
    plot.overlap = -27   # Adjust series overlap to -27%

    # Add data labels with formatted currency
    data_lists = [ast_savings_list, current_savings_list]
    for idx, series in enumerate(chart.series):
        values = data_lists[idx]
        for point_idx, point in enumerate(series.points):
            data_label = point.data_label
            data_label.has_value = False  # We will set the text manually
            data_label.number_format_is_linked = False
            data_label.visible = True

            # Set font properties
            text_frame = data_label.text_frame
            text_frame.text = ''  # Clear existing text
            p = text_frame.paragraphs[0]
            run = p.add_run()
            value = values[point_idx]
            run.text = format_currency_in_millions(value)
            run.font.size = Pt(12)
            run.font.name = "Calibri"
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black

def process_scenario_dataframe(df, scenario_detail_grouping, require_grouping=True):
    """
    Processes the scenario DataFrame to ensure required columns are present.
    If require_grouping=False, we do not insist on the scenario_detail_grouping column.

    Parameters:
    - df: DataFrame to process.
    - scenario_detail_grouping: The grouping column name used for scenario details.
    - require_grouping: boolean, if True scenario_detail_grouping must be included, otherwise it's optional.

    Returns:
    - df: Processed DataFrame with required columns added if necessary.
    """

    # Ensure columns are stripped
    df.columns = df.columns.str.strip()

    # Map existing columns to required columns
    df = df.rename(columns={
        'Awarded Supplier': 'Awarded Supplier Name'
        # Add more mappings if necessary
    })

    # Base expected columns
    expected_columns = [
        'Awarded Supplier Name',
        'Awarded Supplier Spend',
        'AST Savings',
        'Current Savings',
        'AST Baseline Spend',
        'Current Baseline Spend'
    ]

    # Only include scenario_detail_grouping if require_grouping is True
    if require_grouping and scenario_detail_grouping:
        expected_columns.append(scenario_detail_grouping)

    # Calculate 'Awarded Supplier Spend' if not present
    if 'Awarded Supplier Spend' not in df.columns:
        if 'Awarded Supplier Price' in df.columns and 'Awarded Volume' in df.columns:
            df['Awarded Supplier Spend'] = df['Awarded Supplier Price'] * df['Awarded Volume']
        else:
            df['Awarded Supplier Spend'] = 0

    # 'AST Baseline Spend'
    if 'AST Baseline Spend' not in df.columns:
        if 'Baseline Spend' in df.columns:
            df['AST Baseline Spend'] = df['Baseline Spend']
        elif 'Baseline Price' in df.columns and 'Bid Volume' in df.columns:
            df['AST Baseline Spend'] = df['Baseline Price'] * df['Bid Volume']
        else:
            df['AST Baseline Spend'] = 0

    # 'AST Savings'
    if 'AST Savings' not in df.columns:
        df['AST Savings'] = df['AST Baseline Spend'] - df['Awarded Supplier Spend']

    # 'Current Baseline Spend'
    if 'Current Baseline Spend' not in df.columns:
        if 'Current Price' in df.columns and 'Bid Volume' in df.columns:
            df['Current Baseline Spend'] = df['Current Price'] * df['Bid Volume']
        else:
            df['Current Baseline Spend'] = df['AST Baseline Spend']

    # 'Current Savings'
    if 'Current Savings' not in df.columns:
        df['Current Savings'] = df['Current Baseline Spend'] - df['Awarded Supplier Spend']

    # Ensure 'Awarded Volume' exists
    if 'Awarded Volume' not in df.columns:
        if 'Bid Volume' in df.columns:
            df['Awarded Volume'] = df['Bid Volume']
        else:
            df['Awarded Volume'] = 0

    # Ensure 'Incumbent' exists
    if 'Incumbent' not in df.columns:
        df['Incumbent'] = 'Unknown'

    # Ensure 'Bid ID' exists
    if 'Bid ID' not in df.columns:
        df['Bid ID'] = df.index

    # Check expected columns
    for col in expected_columns:
        if col not in df.columns:
            if col == scenario_detail_grouping and require_grouping:
                # If grouping is required but not found, raise error
                raise ValueError(f"Required grouping column '{scenario_detail_grouping}' not found in data.")
            else:
                # For non-grouping columns, just fill with 0 if not found
                df[col] = 0

    return df

def create_scenario_summary_slides(prs, scenario_dataframes, scenario_detail_grouping, title_suffix="", create_details=True, logger=logger):
    """
    Create scenario summary slides for the given scenario_dataframes and grouping.
    If create_details=True, scenario detail slides are created (requiring scenario_detail_grouping).
    If create_details=False, no detail slides are created and scenario_detail_grouping is not required.

    The title_suffix (e.g., " (Port Hudson)") will appear in the header of the slides if provided.
    """
    scenario_keys = list(scenario_dataframes.keys())
    scenarios = [sheet_name.lstrip('#') for sheet_name in scenario_keys]

    logger.info(f"Creating scenario summary slides with title_suffix='{title_suffix}'. Scenarios: {scenarios}")

    scenarios_per_slide = 3
    total_slides = (len(scenarios) + scenarios_per_slide - 1) // scenarios_per_slide

    scenario_index = 0

    for slide_num in range(1, total_slides + 1):
        if slide_num == 1 and len(prs.slides) > 0 and title_suffix == "":
            # Modify the existing first slide only if this is the main scenario summary and no suffix
            slide = prs.slides[0]
            logger.info("Modifying the existing first slide for main scenario summary.")
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
        else:
            default_slide_layout_index = 1
            slide_layout = prs.slide_layouts[default_slide_layout_index]
            slide = prs.slides.add_slide(slide_layout)
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

        add_header(slide, slide_num, title_suffix)
        add_row_labels(slide)

        scenario_names = []
        ast_savings_list = []
        current_savings_list = []

        for i in range(scenarios_per_slide):
            if scenario_index >= len(scenarios):
                break
            scenario_name = scenarios[scenario_index]
            scenario_key = scenario_keys[scenario_index]
            df = scenario_dataframes[scenario_key]

            # Process df with or without requiring grouping
            df = process_scenario_dataframe(df, scenario_detail_grouping, require_grouping=create_details)
            df = df.fillna(0).replace([float('inf'), float('-inf')], 0)

            add_scenario_content(slide, df, scenario_name, scenario_position=i+1)

            total_ast_savings = df['AST Savings'].sum()
            total_current_savings = df['Current Savings'].sum()
            if pd.isna(total_ast_savings):
                total_ast_savings = 0
            if pd.isna(total_current_savings):
                total_current_savings = 0

            ast_savings_list.append(total_ast_savings)
            current_savings_list.append(total_current_savings)
            scenario_names.append(scenario_name)

            if create_details:
                detail_slide_layout_index = 1
                add_scenario_detail_slide(prs, df, scenario_name, detail_slide_layout_index, scenario_detail_grouping)

            scenario_index += 1

        add_chart(slide, scenario_names, ast_savings_list, current_savings_list)

    return prs

def create_scenario_summary_presentation(scenario_dataframes, template_file_path=None):
    """
    Creates a Presentation object based on the scenario DataFrames and a template file,
    including main scenario summaries and optional sub-summaries if toggled on.
    """
    try:
        if template_file_path and os.path.exists(template_file_path):
            prs = Presentation(template_file_path)
            logger.info(f"Loaded PowerPoint template from {template_file_path}")
        else:
            prs = Presentation()
            if template_file_path:
                logger.warning("Template file not found. Using default presentation.")
            else:
                logger.info("No template file provided. Using default presentation.")

        scenario_detail_grouping = st.session_state.get('scenario_detail_grouping', None)
        on = st.session_state.get("scenario_sub_summaries_on", False)
        scenario_summary_selections = st.session_state.get("scenario_summary_selections", None)
        sub_summaries_list = st.session_state.get("sub_summary_selections", [])

        logger.info(f"Sub-summaries on: {on}")
        logger.info(f"Scenario summary selections column: {scenario_summary_selections}")
        logger.info(f"Chosen sub-summaries: {sub_summaries_list}")

        # Create main scenario summaries with details (requires scenario_detail_grouping)
        prs = create_scenario_summary_slides(
            prs=prs,
            scenario_dataframes=scenario_dataframes,
            scenario_detail_grouping=scenario_detail_grouping,
            title_suffix="",
            create_details=True,
            logger=logger
        )

        # Handle sub-summaries if toggled on
        if on and scenario_summary_selections and sub_summaries_list:
            logger.info(f"Creating sub-summary slides for {scenario_summary_selections} values: {sub_summaries_list}")
            for sub_val in sub_summaries_list:
                filtered_dataframes = {}
                for sheet_name, df in scenario_dataframes.items():
                    if scenario_summary_selections in df.columns:
                        sub_filtered_df = df[df[scenario_summary_selections] == sub_val]
                        if not sub_filtered_df.empty:
                            sub_filtered_df = sub_filtered_df.fillna(0).replace([float('inf'), float('-inf')], 0)
                            filtered_dataframes[sheet_name] = sub_filtered_df.copy()

                if filtered_dataframes:
                    logger.info(f"Creating sub-summary slides for '{sub_val}'.")
                    # For sub-summaries, we only show scenario summaries (no details)
                    # Thus create_details=False and no requirement for scenario_detail_grouping
                    prs = create_scenario_summary_slides(
                        prs=prs,
                        scenario_dataframes=filtered_dataframes,
                        scenario_detail_grouping=None,  # Not required here
                        title_suffix=f" ({sub_val})",
                        create_details=False,
                        logger=logger
                    )
                else:
                    logger.warning(f"No data found for sub-summary '{sub_val}'. No additional slides created.")

        return prs

    except Exception as e:
        st.error(f"An error occurred while generating the presentation: {str(e)}")
        logger.error(f"Error generating presentation: {str(e)}")
        return None



def add_scenario_detail_slide(prs, df, scenario_name, template_slide_layout_index, scenario_detail_grouping):
    """
    Adds a detailed slide for the given scenario to the presentation.

    Parameters:
    - prs: The PowerPoint presentation object.
    - df: The DataFrame containing the scenario data, including the grouping column.
    - scenario_name: The name of the scenario.
    - template_slide_layout_index: The index of the slide layout to use.
    - scenario_detail_grouping: The column name to group data by.
    """
    # Necessary imports


    # Ensure the grouping column is in the DataFrame
    if scenario_detail_grouping not in df.columns:
        print(f"The selected grouping field '{scenario_detail_grouping}' is not present in the data.")
        return

    # Add a new slide using the specified layout
    slide_layout = prs.slide_layouts[template_slide_layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # Remove existing shapes if necessary
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Add Header
    left = Inches(0)
    top = Inches(0.01)
    width = Inches(12.5)
    height = Inches(0.5)
    header = slide.shapes.add_textbox(left, top, width, height)
    header_tf = header.text_frame
    header_tf.vertical_anchor = MSO_ANCHOR.TOP
    header_tf.word_wrap = True
    p = header_tf.paragraphs[0]
    p.text = f"{scenario_name} Scenario Details"
    p.font.size = Pt(25)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0, 51, 153)  # Dark Blue
    p.font.bold = True
    header.fill.background()  # No Fill
    header.line.fill.background()  # No Line

    # Calculate percentage of awarded volume for each supplier
    total_awarded_volume = df['Awarded Volume'].sum()
    supplier_volumes = df.groupby('Awarded Supplier Name')['Awarded Volume'].sum().reset_index()
    supplier_volumes['Volume %'] = supplier_volumes['Awarded Volume'] / total_awarded_volume

    # Sort suppliers by 'Volume %' in descending order
    supplier_volumes = supplier_volumes.sort_values(by='Volume %', ascending=False)

    # Prepare data for the chart
    suppliers = supplier_volumes['Awarded Supplier Name'].tolist()
    volume_percentages = supplier_volumes['Volume %'].tolist()

    # Create a consistent color mapping for suppliers
    unique_suppliers = list(OrderedDict.fromkeys(df['Awarded Supplier Name'].tolist() + df['Incumbent'].tolist()))
    # Expanded color palette to include 25 colors
    colorful_palette_3 = [
        RGBColor(68, 114, 196),   # Dark Blue
        RGBColor(237, 125, 49),   # Orange
        RGBColor(165, 165, 165),  # Gray
        RGBColor(255, 192, 0),    # Gold
        RGBColor(112, 173, 71),   # Green
        RGBColor(91, 155, 213),   # Light Blue
        RGBColor(193, 152, 89),   # Brown
        RGBColor(155, 187, 89),   # Olive Green
        RGBColor(128, 100, 162),  # Purple
        RGBColor(158, 72, 14),    # Dark Orange
        RGBColor(99, 99, 99),     # Dark Gray
        RGBColor(133, 133, 133),  # Medium Gray
        RGBColor(49, 133, 156),   # Teal
        RGBColor(157, 195, 230),  # Sky Blue
        RGBColor(75, 172, 198),   # Aqua
        RGBColor(247, 150, 70),   # Light Orange
        RGBColor(128, 128, 0),    # Olive
        RGBColor(192, 80, 77),    # Dark Red
        RGBColor(0, 176, 80),     # Bright Green
        RGBColor(79, 129, 189),   # Steel Blue
        RGBColor(192, 0, 0),      # Red
        RGBColor(0, 112, 192),    # Medium Blue
        RGBColor(0, 176, 240),    # Cyan
        RGBColor(255, 0, 0),      # Bright Red
        RGBColor(146, 208, 80),   # Light Green
    ]
    color_cycle = cycle(colorful_palette_3)
    supplier_color_map = {}
    for supplier in unique_suppliers:
        supplier_color_map[supplier] = next(color_cycle)

    # Create the horizontal stacked bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['']

    for supplier, volume_pct in zip(suppliers, volume_percentages):
        chart_data.add_series(supplier, [volume_pct])

    # Define the chart size and position
    x, y, cx, cy = Inches(0), Inches(0.3), Inches(9.0), Inches(1.7)  # Adjusted width to 9.0 inches
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
    ).chart

    # Format the chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "% of Volume"

    # Ensure legend is created
    chart.has_legend = True
    if chart.legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = "Calibri"

    # Hide the value axis (vertical axis)
    if chart.value_axis:
        chart.value_axis.visible = False

    # Format data labels and series colors
    for idx, series in enumerate(chart.series):
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_value = True
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.INSIDE_BASE
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = RGBColor(255, 255, 255)  # White

        # Set series color
        supplier = series.name
        color = supplier_color_map.get(supplier, RGBColor(0x00, 0x00, 0x00))  # Default to black
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Calculate Transitions and Items
    num_bid_ids = df['Bid ID'].nunique()
    transitions_df = df[df['Awarded Supplier Name'] != df['Incumbent']]
    num_transitions = transitions_df['Bid ID'].nunique()

    # Add Transitions Box
    left = Inches(11.14)
    top = Inches(0.65)
    width = Inches(2.0)
    height = Inches(1.0)
    transitions_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    transitions_box.fill.solid()
    transitions_box.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    transitions_box.line.fill.background()  # No Line

    # Add text to the transitions box
    tb = transitions_box.text_frame
    tb.text = f"# of Transitions\n{num_transitions}"
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Set alignment and font properties for all paragraphs
    for paragraph in tb.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White

    # Add Items Box
    left = Inches(8.97)
    items_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    items_box.fill.solid()
    items_box.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
    items_box.line.fill.background()  # No Line

    # Add text to the items box
    tb = items_box.text_frame
    tb.text = f"# of items\n{num_bid_ids}"
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Set alignment and font properties for all paragraphs
    for paragraph in tb.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White

    # Process data for the table
    grouped_df = df.groupby(scenario_detail_grouping)
    summary_data = []

    for group_name, group_df in grouped_df:
        bid_volume = group_df['Bid Volume'].sum()
        avg_current_price = group_df['Current Price'].mean()
        current_spend = avg_current_price * bid_volume
        avg_bid_price = group_df['Awarded Supplier Price'].mean()
        bid_spend = avg_bid_price * bid_volume
        current_savings = group_df['Current Savings'].sum()

        # Collect Incumbent and Awarded Supplier distributions
        incumbent_dist = group_df.groupby('Incumbent')['Bid Volume'].sum().reset_index()
        awarded_supplier_dist = group_df.groupby('Awarded Supplier Name')['Bid Volume'].sum().reset_index()

        # Append data to the summary list
        summary_data.append({
            'Grouping': group_name,
            'Bid Volume': bid_volume,
            'Incumbent Dist': incumbent_dist,
            'Avg Current Price': avg_current_price,
            'Current Spend': current_spend,
            'Awarded Supplier Dist': awarded_supplier_dist,
            'Avg Bid Price': avg_bid_price,
            'Bid Spend': bid_spend,
            'Current Savings': current_savings
        })

    summary_df = pd.DataFrame(summary_data)

    # Create the table
    rows = len(summary_df) + 2  # Including header row and totals row
    cols = 9  # Number of columns

    left = Inches(0.5)
    top = Inches(2.0)  # Table starts 2 inches from the top
    width = Inches(12.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, Inches(1)).table  # Initial height

    # Set column headers
    column_headers = [
        scenario_detail_grouping,
        'Bid Volume',
        'Incumbent',
        'Avg Current Price',
        'Current Spend',
        'Awarded Supplier',
        'Avg Bid Price',
        'Bid Spend',
        'Current Savings'
    ]

    for col_idx, header in enumerate(column_headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        # Apply header formatting
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.name = "Calibri"
        paragraph.alignment = PP_ALIGN.CENTER

    # Adjust column widths to accommodate pie charts
    table.columns[2].width = Inches(1.5)  # Incumbent
    table.columns[5].width = Inches(1.5)  # Awarded Supplier

    # Set header row height
    table.rows[0].height = Inches(0.38)

    # Set data rows heights
    for row_idx in range(1, len(summary_df) + 1):
        table.rows[row_idx].height = Inches(1.0)

    # Set totals row height
    total_row_idx = len(summary_df) + 1
    table.rows[total_row_idx].height = Inches(0.38)

    # Populate table rows and add pie charts
    for row_idx, data in enumerate(summary_data, start=1):
        table.cell(row_idx, 0).text = str(data['Grouping'])
        table.cell(row_idx, 1).text = f"{data['Bid Volume']:,.0f}"
        table.cell(row_idx, 3).text = f"${data['Avg Current Price']:.2f}"
        table.cell(row_idx, 4).text = format_currency_in_millions(data['Current Spend'])
        table.cell(row_idx, 6).text = f"${data['Avg Bid Price']:.2f}"
        table.cell(row_idx, 7).text = format_currency_in_millions(data['Bid Spend'])
        table.cell(row_idx, 8).text = format_currency_in_millions(data['Current Savings'])

        # Center align text in cells
        for col_idx in [0, 1, 3, 4, 6, 7, 8]:
            cell = table.cell(row_idx, col_idx)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(10)
            paragraph.font.name = 'Calibri'

        # Calculate positions for pie charts
        # Incumbent Pie Chart in 'Incumbent' cell (col_idx = 2)
        col_idx = 2
        # Get the position and size of the cell
        cell_left = left + sum([table.columns[i].width for i in range(col_idx)])
        cell_top = top + Inches(0.38) + Inches(1.0) * (row_idx - 1)
        cell_width = table.columns[col_idx].width
        cell_height = table.rows[row_idx].height

        # Create the pie chart data
        incumbent_chart_data = ChartData()
        incumbent_chart_data.categories = data['Incumbent Dist']['Incumbent']
        incumbent_chart_data.add_series('', data['Incumbent Dist']['Bid Volume'])

        # Add the pie chart
        incumbent_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, cell_left, cell_top, cell_width, cell_height, incumbent_chart_data
        ).chart

        # Disable chart title and legend
        incumbent_chart.has_title = False
        incumbent_chart.has_legend = False

        # Set colors for data points
        for idx_point, point in enumerate(incumbent_chart.series[0].points):
            category = data['Incumbent Dist']['Incumbent'].iloc[idx_point]
            color = supplier_color_map.get(category, RGBColor(0x00, 0x00, 0x00))
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color

        # Disable data labels
        for series in incumbent_chart.series:
            series.has_data_labels = False

        # Awarded Supplier Pie Chart in 'Awarded Supplier' cell (col_idx = 5)
        col_idx = 5
        # Get the position and size of the cell
        cell_left = left + sum([table.columns[i].width for i in range(col_idx)])
        cell_top = top + Inches(0.38) + Inches(1.0) * (row_idx - 1)
        cell_width = table.columns[col_idx].width
        cell_height = table.rows[row_idx].height

        # Create the pie chart data
        awarded_chart_data = ChartData()
        awarded_chart_data.categories = data['Awarded Supplier Dist']['Awarded Supplier Name']
        awarded_chart_data.add_series('', data['Awarded Supplier Dist']['Bid Volume'])

        # Add the pie chart
        awarded_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, cell_left, cell_top, cell_width, cell_height, awarded_chart_data
        ).chart

        # Disable chart title and legend
        awarded_chart.has_title = False
        awarded_chart.has_legend = False

        # Set colors for data points
        for idx_point, point in enumerate(awarded_chart.series[0].points):
            category = data['Awarded Supplier Dist']['Awarded Supplier Name'].iloc[idx_point]
            color = supplier_color_map.get(category, RGBColor(0x00, 0x00, 0x00))
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color

        # Disable data labels
        for series in awarded_chart.series:
            series.has_data_labels = False

    # Add Totals row
    total_bid_volume = summary_df['Bid Volume'].sum()
    total_current_spend = summary_df['Current Spend'].sum()
    total_bid_spend = summary_df['Bid Spend'].sum()
    total_current_savings = summary_df['Current Savings'].sum()
    
    # Calculate weighted average prices
    total_avg_current_price = (
        (summary_df['Avg Current Price'] * summary_df['Bid Volume']).sum() / total_bid_volume
    )
    total_avg_bid_price = (
        (summary_df['Avg Bid Price'] * summary_df['Bid Volume']).sum() / total_bid_volume
    )

    total_row_idx = len(summary_df) + 1  # Totals row index

    table.cell(total_row_idx, 0).text = "Totals"
    table.cell(total_row_idx, 1).text = f"{total_bid_volume:,.0f}"
    table.cell(total_row_idx, 3).text = f"${total_avg_current_price:.2f}"
    table.cell(total_row_idx, 4).text = format_currency_in_millions(total_current_spend)
    table.cell(total_row_idx, 6).text = f"${total_avg_bid_price:.2f}"
    table.cell(total_row_idx, 7).text = format_currency_in_millions(total_bid_spend)
    table.cell(total_row_idx, 8).text = format_currency_in_millions(total_current_savings)

    # Bold and center align the totals row
    for col_idx in range(cols):
        cell = table.cell(total_row_idx, col_idx)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Calibri'


# ///// Bid Coverage Summary Helper Functions /////
def create_bid_coverage_summary_slides(prs, df, bid_coverage_slides_grouping):
    # ----------------- Initial Setup -----------------
    column_mapping = st.session_state.column_mapping

    # Strip leading/trailing spaces from all column names
    df.columns = df.columns.str.strip()

    # Column mappings
    bid_id_col = column_mapping.get('Bid ID')
    incumbent_col = column_mapping.get('Incumbent')
    supplier_name_col = column_mapping.get('Supplier Name')
    bid_price_col = column_mapping.get('Bid Price')
    bid_volume_col = column_mapping.get('Bid Volume')
    baseline_price_col = column_mapping.get('Baseline Price')
    supplier_capacity_col = column_mapping.get('Supplier Capacity')
    awarded_supplier_col = column_mapping.get('Awarded Supplier', 'Awarded Supplier')
    current_price_col = column_mapping.get('Current Price', None)

    # Validate column mappings
    required_mappings = ['Bid ID', 'Incumbent', 'Supplier Name', 'Bid Price', 'Bid Volume', 'Baseline Price', 'Supplier Capacity']
    missing_mappings = [key for key in required_mappings if column_mapping.get(key) not in df.columns]
    if missing_mappings:
        st.error(f"The following required columns are missing in the data: {', '.join(missing_mappings)}.")
        return prs

    # Ensure 'Awarded Supplier' exists using column mapping
    if awarded_supplier_col not in df.columns:
        df[awarded_supplier_col] = df[supplier_name_col]

    # Ensure Bid ID is string type
    df[bid_id_col] = df[bid_id_col].astype(str)

    # Retrieve and sort suppliers
    suppliers = df[supplier_name_col].dropna().unique().tolist()
    suppliers = sorted(suppliers)
    if not suppliers:
        st.error("No suppliers found in the data to generate Bid Coverage Summary slides.")
        return prs

    # Retrieve and sort group values
    group_values = df[bid_coverage_slides_grouping].dropna().unique().tolist()
    group_values = sorted(group_values)
    if not group_values:
        st.error(f"No valid group values found for '{bid_coverage_slides_grouping}'.")
        return prs

    all_items_label = "All Items"

    # ----------------- Slide 1: Bid Coverage by Grouping Type -----------------
    def calc_supplier_stats(subset, column_mapping):
        total_bids = subset[column_mapping['Bid ID']].nunique()
        supplier_stats = {}
        for sup in suppliers:
            sup_data = subset[subset[column_mapping['Supplier Name']] == sup]
            sup_bids = sup_data[column_mapping['Bid ID']].nunique()
            coverage_pct = (sup_bids / total_bids * 100) if total_bids > 0 else 0.0
            
            # Retrieve mapped column names
            bid_price = column_mapping['Bid Price']
            bid_volume = column_mapping['Bid Volume']
            current_price = column_mapping.get('Current Price', None)
            
            # Define required columns based on the availability of 'Current Price'
            if current_price and current_price in sup_data.columns:
                required_cols = [current_price, bid_price, bid_volume]
            else:
                required_cols = [bid_price, bid_volume]
            
            # Check if all required columns exist in sup_data
            if all(col in sup_data.columns for col in required_cols):
                sup_data_clean = sup_data.dropna(subset=required_cols)
                if current_price:
                    savings = ((sup_data_clean[current_price] - sup_data_clean[bid_price]) * sup_data_clean[bid_volume]).sum()
                else:
                    savings = 0
            else:
                savings = 0
            
            supplier_stats[sup] = {
                'bids': sup_bids,
                'coverage_pct': coverage_pct,
                'savings': savings
            }
        
        return total_bids, supplier_stats

    # Calculate stats for All Items
    all_total_bids, all_supplier_stats = calc_supplier_stats(df, column_mapping)
    
    # Calculate stats for each group
    group_supplier_results = {}
    for g in group_values:
        g_df = df[df[bid_coverage_slides_grouping] == g]
        g_total_bids, g_supplier_stats = calc_supplier_stats(g_df, column_mapping)
        group_supplier_results[g] = (g_total_bids, g_supplier_stats)
    
    # Create Slide
    slide_layout1 = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(slide_layout1)
    
    # Remove all existing shapes from the slide
    for shape in list(slide1.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Slide 1 Title
    slide_title1 = slide1.shapes.add_textbox(Inches(0), Inches(0.01), Inches(12), Inches(0.42))
    title_frame1 = slide_title1.text_frame
    title_frame1.text = f"Bid Coverage by {bid_coverage_slides_grouping}"
    title_paragraph1 = title_frame1.paragraphs[0]
    title_paragraph1.font.size = Pt(25)
    title_paragraph1.font.name = "Calibri"
    title_paragraph1.font.color.rgb = RGBColor(0, 51, 153)
    title_paragraph1.font.bold = True
    title_paragraph1.alignment = PP_ALIGN.LEFT
    
    num_groups = len(group_values)
    total_rows1 = 1 + 3 * (num_groups + 1)
    cols1 = 2 + len(suppliers)
    left1 = Inches(0.5)
    top1 = Inches(1.2)
    width1 = Inches(12.0)
    height1 = Inches(0.0)
    table1 = slide1.shapes.add_table(total_rows1, cols1, left1, top1, width1, height1).table
    
    table1.cell(0, 0).text = bid_coverage_slides_grouping
    table1.cell(0, 1).text = "# of Bids"
    if suppliers:
        for i, sup in enumerate(suppliers, start=2):
            table1.cell(0, i).text = sup
    
    for c in range(cols1):
        p = table1.cell(0, c).text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.CENTER
    
    def fill_group_block(start_row, group_label, total_bids, supplier_stats):
        table1.cell(start_row, 0).text = group_label
        table1.cell(start_row, 1).text = f"{total_bids:,.0f}"
        if suppliers:
            for i, sup in enumerate(suppliers, start=2):
                sbids = supplier_stats[sup]['bids']
                table1.cell(start_row, i).text = f"{sbids:,.0f}"
    
        table1.cell(start_row + 1, 0).text = "Item Coverage %"
        table1.cell(start_row + 1, 1).text = ""
        if suppliers:
            for i, sup in enumerate(suppliers, start=2):
                c_pct = supplier_stats[sup]['coverage_pct']
                table1.cell(start_row + 1, i).text = f"{c_pct:.0f}%"
    
        table1.cell(start_row + 2, 0).text = "Savings"
        table1.cell(start_row + 2, 1).text = ""
        if suppliers:
            for i, sup in enumerate(suppliers, start=2):
                savings = supplier_stats[sup]['savings']
                table1.cell(start_row + 2, i).text = format_currency_in_millions(savings)
    
                # Conditional formatting for savings row
                p = table1.cell(start_row + 2, i).text_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.bold = True
                if savings < 0:
                    p.font.color.rgb = RGBColor(255, 0, 0)  # Red for negative
                elif savings > 0:
                    p.font.color.rgb = RGBColor(0, 128, 0)  # Green for positive
                else:
                    p.font.color.rgb = RGBColor(0, 0, 0)   # Black for zero
    
        # Formatting for all rows in the block
        for r in range(start_row, start_row + 3):
            for c in range(cols1):
                cell_par = table1.cell(r, c).text_frame.paragraphs[0]
                cell_par.font.size = Pt(10)
                if c == 0:
                    cell_par.alignment = PP_ALIGN.LEFT
                    cell_par.font.bold = True
                else:
                    # If it's not the savings row (start_row+2), no bold unless savings row handled above
                    if r != (start_row + 2):
                        cell_par.font.bold = False
                    cell_par.alignment = PP_ALIGN.CENTER
    
    # Fill All Items block
    row_index1 = 1
    fill_group_block(row_index1, all_items_label, all_total_bids, all_supplier_stats)
    row_index1 += 3
    
    # Fill each group block
    for g in group_values:
        g_total_bids, g_supplier_stats = group_supplier_results[g]
        fill_group_block(row_index1, str(g), g_total_bids, g_supplier_stats)
        row_index1 += 3

    # ----------------- Slide 2: Bid Coverage by Grouping Type Tables -----------------
    def calc_metrics_slide2(subset):
        total_items = subset[bid_id_col].nunique()
        new_and_incumbent = 0
        only_new_suppliers = 0
        only_incumbent = 0
        not_covered = 0

        bid_ids = subset[bid_id_col].unique()

        for bid_id in bid_ids:
            bid_df = subset[subset[bid_id_col] == bid_id]
            suppliers_who_bid = bid_df[supplier_name_col].dropna().unique()
            if len(suppliers_who_bid) == 0:
                not_covered += 1
                continue

            incumbent_suppliers = bid_df[incumbent_col].dropna().unique()
            set_suppliers_who_bid = set(suppliers_who_bid)
            set_incumbent_suppliers = set(incumbent_suppliers)

            if set_suppliers_who_bid == set_incumbent_suppliers and len(set_incumbent_suppliers) > 0:
                only_incumbent += 1
            elif set_suppliers_who_bid.isdisjoint(set_incumbent_suppliers):
                only_new_suppliers += 1
            elif set_suppliers_who_bid & set_incumbent_suppliers:
                new_and_incumbent += 1

        return {
            'total_items': total_items,
            'new_and_incumbent': new_and_incumbent,
            'only_new_suppliers': only_new_suppliers,
            'only_incumbent': only_incumbent,
            'not_covered': not_covered
        }

    all_metrics_slide2 = calc_metrics_slide2(df)
    group_results_slide2 = {}
    for g in group_values:
        g_df = df[df[bid_coverage_slides_grouping] == g]
        metrics = calc_metrics_slide2(g_df)
        group_results_slide2[g] = metrics

    table_data_slide2 = []
    for group in group_values:
        metrics = group_results_slide2[group]
        table_data_slide2.append({
            'Group': group,
            '# of Items': metrics['total_items'],
            'New & Incumbent': metrics['new_and_incumbent'],
            'Only New Suppliers': metrics['only_new_suppliers'],
            'Only Incumbent': metrics['only_incumbent'],
            'Not Covered': metrics['not_covered']
        })

    slide_layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout2)
    for shape in list(slide2.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Slide 2 Title
    slide_title2 = slide2.shapes.add_textbox(Inches(0), Inches(0.01), Inches(12), Inches(0.42))
    title_frame2 = slide_title2.text_frame
    title_frame2.text = f"Bid Coverage by {bid_coverage_slides_grouping} Tables"
    title_paragraph2 = title_frame2.paragraphs[0]
    title_paragraph2.font.size = Pt(25)
    title_paragraph2.font.name = "Calibri"
    title_paragraph2.font.color.rgb = RGBColor(0, 51, 153)
    title_paragraph2.font.bold = True
    title_paragraph2.alignment = PP_ALIGN.LEFT

    rows2 = len(table_data_slide2) + 1
    cols2 = 6
    left2 = Inches(0.5)
    top2 = Inches(4.0)
    width2 = Inches(6.0)
    height2 = Inches(3.0)
    table2 = slide2.shapes.add_table(rows2, cols2, left2, top2, width2, height2).table

    headers2 = ['Group', '# of Items', 'New & Incumbent', 'Only New Suppliers', 'Only Incumbent', 'Not Covered']
    for col_idx, header in enumerate(headers2):
        cell = table2.cell(0, col_idx)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    for row_idx, data in enumerate(table_data_slide2, start=1):
        table2.cell(row_idx, 0).text = str(data['Group'])
        table2.cell(row_idx, 1).text = f"{data['# of Items']}"
        table2.cell(row_idx, 2).text = f"{data['New & Incumbent']}"
        table2.cell(row_idx, 3).text = f"{data['Only New Suppliers']}"
        table2.cell(row_idx, 4).text = f"{data['Only Incumbent']}"
        table2.cell(row_idx, 5).text = f"{data['Not Covered']}"

        for col_idx in range(cols2):
            cell = table2.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = "Calibri"
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER

    chart_data2 = CategoryChartData()
    chart_data2.categories = [data['Group'] for data in table_data_slide2]
    chart_data2.add_series('New & Incumbent', [data['New & Incumbent'] for data in table_data_slide2])
    chart_data2.add_series('Only New Suppliers', [data['Only New Suppliers'] for data in table_data_slide2])
    chart_data2.add_series('Only Incumbent', [data['Only Incumbent'] for data in table_data_slide2])
    chart_data2.add_series('Not Covered', [data['Not Covered'] for data in table_data_slide2])

    x2, y2, cx2, cy2 = Inches(0.5), Inches(1.0), Inches(6.0), Inches(3.0)
    chart2 = slide2.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x2, y2, cx2, cy2, chart_data2).chart

    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Count of Bids by " + bid_coverage_slides_grouping
    title_paragraph2_chart = chart2.chart_title.text_frame.paragraphs[0]
    title_paragraph2_chart.font.size = Pt(9)
    title_paragraph2_chart.font.name = "Calibri"
    title_paragraph2_chart.font.color.rgb = RGBColor(0, 51, 153)
    title_paragraph2_chart.font.bold = True
    title_paragraph2_chart.alignment = PP_ALIGN.LEFT

    category_axis2 = chart2.category_axis
    category_axis2.has_title = True
    category_axis2.axis_title.text_frame.text = bid_coverage_slides_grouping
    cp = category_axis2.axis_title.text_frame.paragraphs[0]
    cp.font.size = Pt(9)
    cp.font.name = "Calibri"
    cp.font.color.rgb = RGBColor(0, 51, 153)
    cp.font.bold = True

    value_axis2 = chart2.value_axis
    value_axis2.has_title = True
    value_axis2.axis_title.text_frame.text = "Count of Bids"
    vp = value_axis2.axis_title.text_frame.paragraphs[0]
    vp.font.size = Pt(9)
    vp.font.name = "Calibri"
    vp.font.color.rgb = RGBColor(0, 51, 153)
    vp.font.bold = True

    chart2.has_legend = True
    if chart2.has_legend and chart2.legend:
        chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart2.legend.include_in_layout = False
        legend_font = chart2.legend.font
        legend_font.size = Pt(8)
        legend_font.name = "Calibri"
        legend_font.color.rgb = RGBColor(0, 51, 153)
        legend_font.bold = True
    else:
        st.warning("Chart legend is not available for Slide 2.")

    try:
        max_value = max([
            data['New & Incumbent'] + data['Only New Suppliers'] + data['Only Incumbent'] + data['Not Covered']
            for data in table_data_slide2
        ]) if table_data_slide2 else 50
        value_axis2.maximum_scale = (int(max_value / 10) + 1) * 10 if max_value > 0 else 50
        value_axis2.minimum_scale = 0
        value_axis2.major_unit = 10
    except Exception as e:
        st.warning(f"Could not set y-axis scale dynamically: {e}")
        value_axis2.maximum_scale = 30
        value_axis2.minimum_scale = 0
        value_axis2.major_unit = 2

    # Box Shapes on Slide 2
    box_width = Inches(2.0)
    box_height = Inches(0.5)
    box_left = Inches(7.0)
    box_tops = [Inches(1.0), Inches(2.5), Inches(4.0)]
    box_texts = ["Subject 1", "Subject 2", "Subject 3"]
    for top, text in zip(box_tops, box_texts):
        box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, top, box_width, box_height)
        box.text = text
        for paragraph in box.text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.name = "Calibri"
            paragraph.font.color.rgb = RGBColor(0, 51, 153)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.LEFT
        box.fill.background()
        line = box.line
        line.color.rgb = RGBColor(0, 51, 153)
        line.width = Pt(1)

    # ----------------- Slide 3: Supplier Coverage - Incumbent vs. New Business -----------------

    max_suppliers_per_table = 4

    totals_slide3 = {'# of Items': 0}
    for supplier in suppliers:
        totals_slide3[supplier] = {'incumbent_count': 0, 'new_count': 0, 'incumbent_total': 0, 'new_total': 0}

    table_data_slide3 = []

    for group in group_values:
        group_df = df[df[bid_coverage_slides_grouping] == group]
        total_items = group_df[bid_id_col].nunique()
        totals_slide3['# of Items'] += total_items

        row_data = {'Group': group, '# of Items': total_items}

        for supplier in suppliers:
            incumbent_total = group_df[group_df[incumbent_col] == supplier][bid_id_col].nunique()
            incumbent_count = group_df[
                (group_df[supplier_name_col] == supplier) &
                (group_df[incumbent_col] == supplier)
            ][bid_id_col].nunique()

            new_total = total_items - incumbent_total
            new_count = group_df[
                (group_df[supplier_name_col] == supplier) &
                (group_df[incumbent_col] != supplier)
            ][bid_id_col].nunique()

            row_data[f"{supplier} Incumbent Count"] = incumbent_count
            row_data[f"{supplier} Incumbent Total"] = incumbent_total
            row_data[f"{supplier} New Count"] = new_count
            row_data[f"{supplier} New Total"] = new_total

            totals_slide3[supplier]['incumbent_count'] += incumbent_count
            totals_slide3[supplier]['incumbent_total'] += incumbent_total
            totals_slide3[supplier]['new_count'] += new_count
            totals_slide3[supplier]['new_total'] += new_total

        table_data_slide3.append(row_data)

    totals_row_slide3 = {'Group': 'Totals', '# of Items': totals_slide3['# of Items']}
    for supplier in suppliers:
        totals_row_slide3[f"{supplier} Incumbent Count"] = totals_slide3[supplier]['incumbent_count']
        totals_row_slide3[f"{supplier} Incumbent Total"] = totals_slide3[supplier]['incumbent_total']
        totals_row_slide3[f"{supplier} New Count"] = totals_slide3[supplier]['new_count']
        totals_row_slide3[f"{supplier} New Total"] = totals_slide3[supplier]['new_total']

    table_data_slide3.append(totals_row_slide3)

    # Create the third slide
    slide_layout3 = prs.slide_layouts[1]  # Assuming a blank layout
    slide3 = prs.slides.add_slide(slide_layout3)

    # Remove existing shapes
    for shape in list(slide3.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Add Slide 3 Title (remain at 25 pt)
    slide_title3 = slide3.shapes.add_textbox(Inches(0), Inches(0.01), Inches(12), Inches(0.42))
    title_frame3 = slide_title3.text_frame
    title_frame3.text = "Supplier Coverage - Incumbent vs. New Business"
    title_paragraph3 = title_frame3.paragraphs[0]
    title_paragraph3.font.size = Pt(25)
    title_paragraph3.font.name = "Calibri"
    title_paragraph3.font.color.rgb = RGBColor(0, 51, 153)
    title_paragraph3.font.bold = True
    title_paragraph3.alignment = PP_ALIGN.LEFT

    from math import ceil
    from pptx.dml.color import RGBColor as dmlRGB
    grey_color = dmlRGB(217, 217, 217)

    # Increase spacing between tables
    table_spacing = Inches(0.5)
    num_chunks = ceil(len(suppliers) / max_suppliers_per_table)

    def create_supplier_table(slide, top, chunk_suppliers, table_data, totals_data):
        data_rows_count = len(table_data)
        rows = 3 + data_rows_count  # 3 header rows + data rows
        cols = 2 + len(chunk_suppliers)*4
        left = Inches(0.5)
        width = Inches(12.0)
        height = Inches(0.0)

        # Add table
        graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = graphic_frame.table
        # Disable autofit
        table.allow_autofit = False

        # Set some column widths to prevent wrapping
        # For simplicity, just divide width roughly
        # 2 columns for grouping + # of items = let's say 1.5 in total width
        # Remaining columns for suppliers:
        # total width = 12 in - 1.5 in = 10.5 in for suppliers
        # Each supplier = 4 columns, each chunk_suppliers * 4 columns
        # For demonstration, just pick a width that seems reasonable
        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(0.5)
        # The rest of the columns:
        remaining_width = 12.0 - 1.5  # 10.5 in
        columns_per_supplier = 4
        total_supplier_cols = len(chunk_suppliers)*columns_per_supplier
        if total_supplier_cols > 0:
            col_width = Inches(10.5 / total_supplier_cols)
            for c in range(2, cols):
                table.columns[c].width = col_width

        # Row 0 headers
        table.cell(0,0).text = bid_coverage_slides_grouping
        table.cell(0,1).text = f"# of Items"

        for idx, supplier in enumerate(chunk_suppliers):
            start_col = 2 + idx*4
            end_col = start_col + 3
            top_left_cell = table.cell(0, start_col)
            top_left_cell.text = supplier
            top_left_cell.merge(table.cell(0, end_col))

        # Row 1 (Incumbent/New)
        table.cell(1,0).text = ""
        table.cell(1,1).text = ""
        for idx, supplier in enumerate(chunk_suppliers):
            start_col = 2 + idx*4
            inc_cell = table.cell(1, start_col)
            inc_cell.text = "Incumbent"
            inc_cell.merge(table.cell(1, start_col+1))
            new_cell = table.cell(1, start_col+2)
            new_cell.text = "New"
            new_cell.merge(table.cell(1, start_col+3))

        # Row 2 (Count/Total)
        table.cell(2,0).text = ""
        table.cell(2,1).text = ""
        for idx, supplier in enumerate(chunk_suppliers):
            start_col = 2 + idx*4
            table.cell(2, start_col).text = "Count"
            table.cell(2, start_col+1).text = "Total"
            table.cell(2, start_col+2).text = "Count"
            table.cell(2, start_col+3).text = "Total"

        data_start_row = 3
        totals_row_index = rows - 1  # The last row is totals row

        # Fill data rows
        for r_idx, data_row in enumerate(table_data, start=data_start_row):
            table.cell(r_idx, 0).text = str(data_row['Group'])
            table.cell(r_idx, 1).text = str(data_row['# of Items'])

            if r_idx < totals_row_index:
                # Normal data rows
                for sup_idx, supplier in enumerate(chunk_suppliers):
                    inc_count_val = data_row.get(f"{supplier} Incumbent Count", 0)
                    inc_total_val = data_row.get(f"{supplier} Incumbent Total", 0)
                    new_count_val = data_row.get(f"{supplier} New Count", 0)
                    new_total_val = data_row.get(f"{supplier} New Total", 0)

                    start_col = 2 + sup_idx * 4
                    table.cell(r_idx, start_col).text = str(inc_count_val)
                    table.cell(r_idx, start_col+1).text = str(inc_total_val)
                    table.cell(r_idx, start_col+2).text = str(new_count_val)
                    table.cell(r_idx, start_col+3).text = str(new_total_val)

        # Highlight total columns for data rows
        for r in range(data_start_row, totals_row_index):
            for sup_idx in range(len(chunk_suppliers)):
                start_col = 2 + sup_idx*4
                it_cell = table.cell(r, start_col+1)  # Incumbent Total
                it_cell.fill.solid()
                it_cell.fill.fore_color.rgb = grey_color
                nt_cell = table.cell(r, start_col+3)  # New Total
                nt_cell.fill.solid()
                nt_cell.fill.fore_color.rgb = grey_color

        # Totals row percentages
        totals_data_row = totals_data
        table.cell(totals_row_index, 0).text = str(totals_data_row['Group'])
        table.cell(totals_row_index, 1).text = str(totals_data_row['# of Items'])

        for sup_idx, supplier in enumerate(chunk_suppliers):
            start_col = 2 + sup_idx*4
            inc_count_val = totals_data_row.get(f"{supplier} Incumbent Count", 0)
            inc_total_val = totals_data_row.get(f"{supplier} Incumbent Total", 0)
            new_count_val = totals_data_row.get(f"{supplier} New Count", 0)
            new_total_val = totals_data_row.get(f"{supplier} New Total", 0)

            inc_pct = 0 if inc_total_val == 0 else (inc_count_val / inc_total_val * 100)
            new_pct = 0 if new_total_val == 0 else (new_count_val / new_total_val * 100)

            # Incumbent percentage cell (merge Count & Total)
            inc_cell = table.cell(totals_row_index, start_col)
            inc_cell.merge(table.cell(totals_row_index, start_col+1))
            inc_cell.text = f"{inc_pct:.0f}%"
            inc_cell.fill.solid()
            inc_cell.fill.fore_color.rgb = grey_color

            # New percentage cell (merge Count & Total)
            new_cell = table.cell(totals_row_index, start_col+2)
            new_cell.merge(table.cell(totals_row_index, start_col+3))
            new_cell.text = f"{new_pct:.0f}%"
            new_cell.fill.solid()
            new_cell.fill.fore_color.rgb = grey_color

        # Set font size to 8 for entire table
        for row_i in range(rows):
            for col_i in range(cols):
                cell = table.cell(row_i, col_i)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = "Calibri"

        # Set a fixed row height so we know exactly how tall the table is
        row_height_pts = 15  # points per row (~0.208 inches)
        for row_obj in table.rows:
            row_obj.height = Pt(row_height_pts)

        # Compute total table height from row count * row_height
        total_table_height_pts = rows * row_height_pts
        total_table_height_in = total_table_height_pts / 72.0

        return top + Inches(total_table_height_in)

    current_top = Inches(1.5)
    for chunk_idx in range(num_chunks):
        chunk_suppliers = suppliers[chunk_idx*max_suppliers_per_table:(chunk_idx+1)*max_suppliers_per_table]
        current_top = create_supplier_table(
            slide3,
            current_top,
            chunk_suppliers,
            table_data_slide3,
            table_data_slide3[-1]
        ) + Inches(0.5)  # increased spacing between tables

    return prs


def create_bid_coverage_summary_presentation(merged_data, bid_coverage_slides_grouping, template_file_path):
    """
    Creates the Bid Coverage Summary presentation.
    This will:
    1. Make a copy of merged_data
    2. Create a ppt presentation with 3 slides.
       Slide 1: Table as specified
       Slide 2 and 3: blank placeholders
    
    Returns: A pptx Presentation object
    """
    from pptx import Presentation

    # Create a copy of merged_data
    df = merged_data.copy()

    # Load template or create blank
    if template_file_path and os.path.exists(template_file_path):
        prs = Presentation(template_file_path)
    else:
        prs = Presentation()

    # Create slides
    prs = create_bid_coverage_summary_slides(prs, df, bid_coverage_slides_grouping)

    return prs



# /// Suppplier Comparison Helper Functions
def create_supplier_comparison_summary_slide(prs, df, grouping):
    """
    Create a slide for the Supplier Comparison Summary.
    This shows a clustered column chart of percentage difference from *current price* as baseline by grouping.
    The difference is computed using Bid Price relative to Current Price.
    """

    # Add a new slide (blank layout)
    slide_layout = prs.slide_layouts[1]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_left = Inches(0.5)
    title_top = Inches(0.2)
    title_width = Inches(9)
    title_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "Supplier Comparison Summary"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    # Ensure grouping column exists
    if grouping not in df.columns:
        return prs

    # Retrieve required mappings from session state
    bid_id_col = st.session_state.column_mapping.get('Bid ID')
    incumbent_col = st.session_state.column_mapping.get('Incumbent')
    supplier_name_col = st.session_state.column_mapping.get('Supplier Name')
    baseline_price_col = st.session_state.column_mapping.get('Baseline Price')  # Still defined, but we'll not use as baseline now
    current_price_col = st.session_state.column_mapping.get('Current Price')    # Using this as our baseline reference
    bid_price_col = st.session_state.column_mapping.get('Bid Price')            # Bid Price used for difference calculation
    bid_volume_col = st.session_state.column_mapping.get('Bid Volume')

    groups = df[grouping].dropna().unique()
    groups = sorted(groups)
    suppliers = df[supplier_name_col].dropna().unique()
    suppliers = sorted(suppliers)

    # Calculate percentage difference from the current price baseline using Bid Price
    def pct_diff_func(x):
        total_vol = x[bid_volume_col].sum()
        if total_vol == 0:
            return 0
        avg_current = (x[current_price_col]*x[bid_volume_col]).sum()/total_vol
        avg_bid = (x[bid_price_col]*x[bid_volume_col]).sum()/total_vol
        if avg_current == 0:
            return 0
        # Return fraction (e.g. 0.1 = 10%)
        return (avg_bid - avg_current)/avg_current

    grouped = df.groupby([grouping, supplier_name_col])
    summary = grouped.apply(pct_diff_func).reset_index(name='pct_diff')

    pivot = summary.pivot(index=grouping, columns=supplier_name_col, values='pct_diff').fillna(0)
    pivot = pivot.reindex(index=groups, columns=suppliers, fill_value=0)

    # Chart location and size
    chart_left = Inches(0)
    chart_top = Inches(1)
    chart_width = Inches(8.3)
    chart_height = Inches(3)

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import MSO_THEME_COLOR

    chart_data = CategoryChartData()
    chart_data.categories = list(pivot.index)
    for supplier in pivot.columns:
        chart_data.add_series(supplier, pivot[supplier].tolist())

    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data)
    chart = chart_shape.chart

    # Chart formatting
    chart.has_title = True
    chart.chart_title.text_frame.text = f"Percentage from Current by {grouping}"
    title_p = chart.chart_title.text_frame.paragraphs[0]
    title_p.font.size = Pt(8)
    title_p.font.name = "Calibri"
    title_p.font.bold = True

    category_axis = chart.category_axis
    category_axis.has_title = True
    category_axis.axis_title.text_frame.text = grouping
    cat_p = category_axis.axis_title.text_frame.paragraphs[0]
    cat_p.font.size = Pt(8)
    cat_p.font.name = "Calibri"
    cat_p.font.bold = True
    category_axis.tick_labels.font.size = Pt(8)
    category_axis.tick_labels.font.name = "Calibri"

    value_axis = chart.value_axis
    value_axis.has_title = True
    value_axis.axis_title.text_frame.text = "% Difference from Current"
    val_p = value_axis.axis_title.text_frame.paragraphs[0]
    val_p.font.size = Pt(8)
    val_p.font.name = "Calibri"
    val_p.font.bold = True

    value_axis.minimum_scale = -0.5
    value_axis.maximum_scale = 0.5
    value_axis.major_unit = 0.1
    value_axis.tick_labels.number_format = '0%'
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.name = "Calibri"

    chart.has_legend = True
    if chart.has_legend and chart.legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        legend_font = chart.legend.font
        legend_font.size = Pt(8)
        legend_font.name = "Calibri"
        legend_font.bold = True

    theme_colors = [MSO_THEME_COLOR.ACCENT_1, MSO_THEME_COLOR.ACCENT_4, MSO_THEME_COLOR.ACCENT_5, MSO_THEME_COLOR.ACCENT_6]
    for i, series in enumerate(chart.series):
        color = theme_colors[i % len(theme_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.theme_color = color
        series.invert_if_negative = False
        series.has_data_labels = False
        dl = series.data_labels
        dl.show_value = True
        dl.show_series_name = False
        dl.number_format = '0%'
        dl.font.size = Pt(8)
        dl.font.name = "Calibri"
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Bullet points in upper right quadrant
    text_box_left = Inches(8.3)
    text_box_top = Inches(1)
    text_box_width = Inches(4)
    text_box_height = Inches(2.83)
    text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p1 = text_frame.add_paragraph()
    p1.text = "Average % from Current by Supplier"
    p1.font.bold = True
    for _ in range(3):
        bullet = text_frame.add_paragraph()
        bullet.text = "Fill in comments here"
        bullet.level = 1

    spacer = text_frame.add_paragraph()
    spacer.text = ""

    p2 = text_frame.add_paragraph()
    p2.text = "Average % from Current by Grouping"
    p2.font.bold = True
    for _ in range(3):
        bullet = text_frame.add_paragraph()
        bullet.text = "Fill in comments here"
        bullet.level = 1

    # Main table (suppliers vs groups) at bottom-left quadrant
    table_data = pivot.T  # suppliers as rows, groups as columns
    rows = len(table_data.index) + 1
    cols = len(table_data.columns) + 1
    table_left = Inches(0.5)
    table_top = Inches(4.5)
    table_width = Inches(4.5)
    table_height = Inches(2)
    graphic_frame = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = graphic_frame.table

    table.cell(0,0).text = ""
    for j, grp in enumerate(table_data.columns, start=1):
        cell = table.cell(0,j)
        cell.text = str(grp)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.name = "Calibri"

    for i, sup in enumerate(table_data.index, start=1):
        cell = table.cell(i,0)
        cell.text = str(sup)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.name = "Calibri"

    for i, sup in enumerate(table_data.index, start=1):
        for j, grp in enumerate(table_data.columns, start=1):
            val = table_data.loc[sup, grp]
            cell = table.cell(i,j)
            cell.text = f"{val*100:.1f}%"
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.name = "Calibri"

    # Average % from Current by Supplier table (top right)
    avg_by_supplier = pivot.mean(axis=0)  # fraction
    sup_rows = len(avg_by_supplier.index) + 2  # 1 title row + 1 header row + data rows
    sup_cols = 2
    sup_table_left = Inches(9.59)
    sup_table_top = Inches(4.49)  # just below bullet box (~4.0)
    sup_table_width = Inches(3)
    sup_table_height = Inches(2.33)
    sup_gf = slide.shapes.add_table(sup_rows, sup_cols, sup_table_left, sup_table_top, sup_table_width, sup_table_height)
    sup_table_obj = sup_gf.table

    # Title row for supplier table
    sup_table_obj.cell(0,0).merge(sup_table_obj.cell(0,1))
    sup_table_obj.cell(0,0).text = "Average % from Current by Supplier"
    p = sup_table_obj.cell(0,0).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Header row
    sup_table_obj.cell(1,0).text = "Supplier"
    p = sup_table_obj.cell(1,0).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"

    sup_table_obj.cell(1,1).text = "Avg %"
    p = sup_table_obj.cell(1,1).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"

    for i, sup in enumerate(avg_by_supplier.index, start=2):
        cell = sup_table_obj.cell(i,0)
        cell.text = sup
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.name = "Calibri"

        val = avg_by_supplier[sup]
        cell = sup_table_obj.cell(i,1)
        cell.text = f"{val*100:.1f}%"
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.name = "Calibri"

    # Average % from Current by Grouping table (bottom right)
    avg_by_group = pivot.mean(axis=1)
    grp_rows = len(avg_by_group.index) + 2
    grp_cols = 2
    grp_table_left = Inches(5.83)
    grp_table_top = Inches(4.5)
    grp_table_width = Inches(3)
    grp_table_height = Inches(1.4)
    grp_gf = slide.shapes.add_table(grp_rows, grp_cols, grp_table_left, grp_table_top, grp_table_width, grp_table_height)
    grp_table_obj = grp_gf.table

    # Title row for grouping table
    grp_table_obj.cell(0,0).merge(grp_table_obj.cell(0,1))
    grp_table_obj.cell(0,0).text = "Average % from Current by Grouping"
    p = grp_table_obj.cell(0,0).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Header row
    grp_table_obj.cell(1,0).text = "Grouping"
    p = grp_table_obj.cell(1,0).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"

    grp_table_obj.cell(1,1).text = "Avg %"
    p = grp_table_obj.cell(1,1).text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(8)
    p.font.name = "Calibri"

    for i, grp_val in enumerate(avg_by_group.index, start=2):
        cell = grp_table_obj.cell(i,0)
        cell.text = str(grp_val)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.name = "Calibri"

        val = avg_by_group[grp_val]
        cell = grp_table_obj.cell(i,1)
        cell.text = f"{val*100:.1f}%"
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.name = "Calibri"

    return prs

